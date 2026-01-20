from __future__ import annotations

import argparse
from pathlib import Path
from pptx import Presentation

from .plan_loader import load_plan
from .ollama_client import OllamaClient
from .llm_inference import infer_edges
from .slide_builder import add_slide_from_plan


DEFAULT_MODEL = "qwen2.5:7b"
DEFAULT_OLLAMA_URL = "http://localhost:11434"


def build_ppt_from_plan(
    plan_path: str | Path,
    out_path: str | Path,
    *,
    model: str = DEFAULT_MODEL,
    ollama_base_url: str = DEFAULT_OLLAMA_URL,
    template_path: str | Path | None = None,
    enable_llm: bool = True,
    enable_arrows: bool = True,
    reserve_figure_space: bool = True,
) -> None:
    """
    Build a PowerPoint (.pptx) from a "plan JSON" (Johnson document-like structure).

    - plan_path: path to the JSON plan (required)
    - out_path: destination pptx path (required)
    - model: Ollama model name (default: qwen2.5:7b)
    - ollama_base_url: where Ollama is running (default: http://localhost:11434)
    - template_path: optional .pptx template to start from
    - enable_llm: if False, no LLM calls are made
    - enable_arrows: if True + enable_llm, infer logical arrows between bullet boxes
    - reserve_figure_space: if True, slides with figure_used reserve space for a figure placeholder
    """
    plan_path = Path(plan_path)
    out_path = Path(out_path)

    slides = load_plan(str(plan_path))

    prs = Presentation(str(template_path)) if template_path else Presentation()
    client = OllamaClient(base_url=ollama_base_url)

    for sp in slides:
        edges = []
        if enable_llm and enable_arrows and sp.bullets:
            edges = infer_edges(
                client,
                model=model,
                title=sp.title,
                bullets=sp.bullets,
                speaker_notes=sp.speaker_notes,
            )

        add_slide_from_plan(
            prs,
            title=sp.title,
            bullets=sp.bullets,
            figure_used=sp.figure_used if reserve_figure_space else None,
            edges=edges,
            speaker_notes=sp.speaker_notes,
        )

    prs.save(str(out_path))


def _cli() -> None:
    ap = argparse.ArgumentParser(
        prog="ppt_from_plan",
        description="Generate a PowerPoint from a plan JSON using python-pptx and optional Ollama inference.",
    )
    ap.add_argument("--plan", required=True, help="Path to plan JSON")
    ap.add_argument("--out", required=True, help="Output .pptx path")
    ap.add_argument("--model", default=DEFAULT_MODEL, help=f"Ollama model (default: {DEFAULT_MODEL})")
    ap.add_argument("--ollama", default=DEFAULT_OLLAMA_URL, help="Ollama base URL (default: http://localhost:11434)")
    ap.add_argument("--template", default=None, help="Optional .pptx template path")
    ap.add_argument("--no-llm", action="store_true", help="Disable LLM calls entirely")
    ap.add_argument("--no-arrows", action="store_true", help="Do not infer/draw arrows")
    ap.add_argument("--no-figure-space", action="store_true", help="Do not reserve space for figure placeholder")

    args = ap.parse_args()

    build_ppt_from_plan(
        plan_path=args.plan,
        out_path=args.out,
        model=args.model,
        ollama_base_url=args.ollama,
        template_path=args.template,
        enable_llm=not args.no_llm,
        enable_arrows=not args.no_arrows,
        reserve_figure_space=not args.no_figure_space,
    )


if __name__ == "__main__":
    _cli()
