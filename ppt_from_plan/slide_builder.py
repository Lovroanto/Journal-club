from __future__ import annotations

from dataclasses import dataclass
from typing import Dict, List, Optional, Tuple

import math

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

EMU_PER_INCH = 914400


def emu_to_inches(x_emu: int) -> float:
    return x_emu / EMU_PER_INCH


def get_slide_size_in(prs) -> tuple[float, float]:
    return emu_to_inches(prs.slide_width), emu_to_inches(prs.slide_height)


@dataclass
class Box:
    x: float
    y: float
    w: float
    h: float


def _add_title(
    slide,
    title: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: int = 32,
) -> None:
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)


def _add_speaker_notes(slide, speaker_notes: str) -> None:
    try:
        notes = slide.notes_slide.notes_text_frame
        notes.clear()
        notes.text = speaker_notes or ""
    except Exception:
        pass


def _add_figure_placeholder(slide, caption: str, x: float, y: float, w: float, h: float) -> None:
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(245, 245, 245)
    shp.line.color.rgb = RGBColor(180, 180, 180)

    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = caption or "Figure"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0, 0, 0)


def _add_bullet_box(slide, text: str, box: Box, font_pt: int = 18):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(box.x), Inches(box.y), Inches(box.w), Inches(box.h),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(250, 250, 250)
    shp.line.color.rgb = RGBColor(120, 120, 120)

    tf = shp.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_pt)
    p.font.color.rgb = RGBColor(0, 0, 0)
    return shp


def _port(box: Box, side: str) -> Tuple[float, float]:
    """Return a connection point (in inches) on a box."""
    if side == "left":
        return (box.x, box.y + box.h / 2.0)
    if side == "right":
        return (box.x + box.w, box.y + box.h / 2.0)
    if side == "top":
        return (box.x + box.w / 2.0, box.y)
    if side == "bottom":
        return (box.x + box.w / 2.0, box.y + box.h)
    return (box.x + box.w / 2.0, box.y + box.h / 2.0)


def _clamp(v: float, lo: float, hi: float) -> float:
    return max(lo, min(hi, v))


def _add_line(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2),
    )
    conn.line.color.rgb = RGBColor(60, 60, 60)
    conn.line.width = Pt(2)


def _add_arrow(
    slide,
    a: Box,
    b: Box,
    *,
    slide_w: float,
    bullets_x: float,
) -> None:
    """
    Required arrow style:

    - Connect LEFT-MIDDLE of source box -> LEFT-MIDDLE of target box.
    - Route "around" using the LEFT margin (space between slide edge and bullets_x).
    - Use an elbow path (3 segments) so it clearly does not pass through content.
    - Put a pointy arrowhead at the target, pointing into the target box.
    """
    # Ports: left-middle only
    sx, sy = _port(a, "left")
    tx, ty = _port(b, "left")

    # Define a left-margin lane (outside the bullets column).
    # bullets_x is the left edge of the bullet boxes column.
    # We place the lane somewhere between the slide left margin and bullets_x.
    # If there is little space, we still keep it as left as possible.
    slide_left = 0.0
    safe_pad = max(0.15, 0.015 * slide_w)  # keep away from absolute edge

    # Desired gutter width (~0.6 in or 6% of slide width, whichever smaller)
    desired_gutter = min(0.60, 0.06 * slide_w)

    # Lane should be to the LEFT of bullet boxes:
    lane_x = bullets_x - desired_gutter

    # Clamp lane_x inside slide, but still left of bullets_x
    lane_x = _clamp(lane_x, slide_left + safe_pad, bullets_x - safe_pad)

    # If bullets_x is too close to the edge and lane_x collapses,
    # force a minimal lane slightly left of the boxes (still within slide).
    if lane_x >= bullets_x - safe_pad:
        lane_x = max(slide_left + safe_pad, bullets_x - 0.20)

    # Elbow path:
    # 1) Horizontal out from source left edge to lane
    # 2) Vertical along lane to target y
    # 3) Horizontal into target left edge
    _add_line(slide, sx, sy, lane_x, sy)
    _add_line(slide, lane_x, sy, lane_x, ty)
    _add_line(slide, lane_x, ty, tx, ty)

    # Arrowhead: at (tx, ty), direction is from lane_x -> tx (should be rightwards)
    dx = tx - lane_x
    dy = 0.0
    angle_deg = math.degrees(math.atan2(dy, dx))  # ~0 if pointing right

    head_size = 0.18  # inches
    tri_left = tx - head_size / 2.0
    tri_top = ty - head_size / 2.0

    tri = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        Inches(tri_left), Inches(tri_top),
        Inches(head_size), Inches(head_size),
    )
    tri.fill.solid()
    tri.fill.fore_color.rgb = RGBColor(60, 60, 60)
    try:
        tri.line.fill.background()
    except Exception:
        tri.line.color.rgb = RGBColor(60, 60, 60)

    # Default triangle points up; rotate to point along +x (into the box)
    tri.rotation = angle_deg + 90.0


def add_slide_from_plan(
    prs: Presentation,
    *,
    title: str,
    bullets: List[str],
    figure_used: Optional[str],
    edges: List[Dict],
    speaker_notes: str = "",
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    slide_w, slide_h = get_slide_size_in(prs)

    # Title sizing heuristic based on length
    title_len = len(title.strip()) if title else 0
    if title_len > 60:
        title_h = 0.22 * slide_h
        title_font = 26
    elif title_len > 40:
        title_h = 0.18 * slide_h
        title_font = 28
    else:
        title_h = 0.14 * slide_h
        title_font = 32

    # Percent-based margins and regions
    margin_x = 0.05 * slide_w
    margin_top = 0.05 * slide_h
    margin_bottom = 0.05 * slide_h
    gap_x = 0.03 * slide_w

    title_y = margin_top
    content_y = title_y + title_h
    content_h = slide_h - content_y - margin_bottom

    left_x = margin_x
    usable_w = slide_w - 2 * margin_x

    _add_title(slide, title, left_x, title_y, usable_w, title_h, font_size=title_font)
    _add_speaker_notes(slide, speaker_notes)

    has_fig = bool(figure_used)
    if has_fig:
        fig_w = 0.40 * usable_w
        bullets_w = usable_w - fig_w - gap_x
        bullets_x = left_x
        fig_x = bullets_x + bullets_w + gap_x
    else:
        bullets_x = left_x
        bullets_w = usable_w
        fig_x = fig_w = 0.0

    if has_fig:
        _add_figure_placeholder(slide, str(figure_used), fig_x, content_y, fig_w, content_h)

    n = len(bullets)
    if n == 0:
        return slide

    # 1 or 2 columns
    cols = 2 if (n >= 6 and bullets_w >= 6.0) else 1
    col_gap = 0.02 * usable_w
    col_w = (bullets_w - col_gap * (cols - 1)) / cols

    rows = (n + cols - 1) // cols
    row_gap = 0.02 * slide_h
    box_h = min(1.1, (content_h - row_gap * (rows - 1)) / max(1, rows))

    boxes: List[Box] = []
    idx = 0
    for r in range(rows):
        for c in range(cols):
            if idx >= n:
                break
            x = bullets_x + c * (col_w + col_gap)
            y = content_y + r * (box_h + row_gap)
            b = Box(x=x, y=y, w=col_w, h=box_h)
            boxes.append(b)
            _add_bullet_box(slide, bullets[idx], b, font_pt=18)
            idx += 1

    # Draw arrows
    for e in edges or []:
        try:
            i = int(e.get("from", -1))
            j = int(e.get("to", -1))
        except Exception:
            continue
        if 0 <= i < len(boxes) and 0 <= j < len(boxes) and i != j:
            _add_arrow(
                slide,
                boxes[i],
                boxes[j],
                slide_w=slide_w,
                bullets_x=bullets_x,
            )

    return slide
